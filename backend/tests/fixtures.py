"""Deterministic per-format fixture generators for the regression suite.

Each function returns ``(filename, bytes)`` built from fixed content so output
is reproducible. Binary formats needing a writer lib are guarded — the test
skips that format if the lib is absent rather than failing the whole suite.
"""

from __future__ import annotations

import io

# ── Text-family (inline, always available) ───────────────────────────────────

TXT = ("notes.txt", b"# Heading\n\nA short paragraph of plain text.\n")
MD = ("doc.md", b"# Title\n\n- one\n- two\n\nBody paragraph.\n")
CSV = ("data.csv", b"name,role,city\nAda,Engineer,London\nAlan,Researcher,Leeds\n")
TSV = ("data.tsv", b"name\trole\nAda\tEngineer\n")
JSON = ("data.json", b'{"team":[{"name":"Ada","role":"Engineer"}]}')
XML = ("data.xml", b"<team><member><name>Ada</name></member></team>")
HTML = (
    "page.html",
    b"<html><body><h1>Title</h1><table><tr><th>A</th><th>B</th></tr>"
    b"<tr><td>1</td><td>2</td></tr></table><p>Hi <a href='http://x'>link</a></p></body></html>",
)

# Minimal hand-crafted PDF (pdfminer-readable), deterministic.
PDF_BYTES = b"""%PDF-1.4
1 0 obj<</Type/Catalog/Pages 2 0 R>>endobj
2 0 obj<</Type/Pages/Kids[3 0 R]/Count 1>>endobj
3 0 obj<</Type/Page/Parent 2 0 R/MediaBox[0 0 300 144]/Contents 4 0 R/Resources<</Font<</F1 5 0 R>>>>>>endobj
4 0 obj<</Length 44>>stream
BT /F1 18 Tf 20 100 Td (Hello MDify Pro) Tj ET
endstream endobj
5 0 obj<</Type/Font/Subtype/Type1/BaseFont/Helvetica>>endobj
xref
0 6
0000000000 65535 f
trailer<</Root 1 0 R/Size 6>>
startxref
0
%%EOF"""
PDF = ("doc.pdf", PDF_BYTES)


# ── Binary formats needing a writer lib ──────────────────────────────────────

def docx_bytes() -> tuple[str, bytes]:
    import docx  # python-docx

    d = docx.Document()
    d.add_heading("Report", level=1)
    d.add_heading("Summary", level=2)
    d.add_paragraph("Revenue grew this quarter.")
    t = d.add_table(rows=2, cols=2)
    t.cell(0, 0).text = "A"
    t.cell(0, 1).text = "B"
    t.cell(1, 0).text = "1"
    t.cell(1, 1).text = "2"
    buf = io.BytesIO()
    d.save(buf)
    return ("report.docx", buf.getvalue())


def pptx_bytes() -> tuple[str, bytes]:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "MDify Pro Deck"
    slide.placeholders[1].text = "Bullet one\nBullet two"
    buf = io.BytesIO()
    prs.save(buf)
    return ("deck.pptx", buf.getvalue())


def xlsx_bytes() -> tuple[str, bytes]:
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.append(["Name", "Score"])
    ws.append(["Ada", 99])
    ws.append(["Alan", 97])
    buf = io.BytesIO()
    wb.save(buf)
    return ("metrics.xlsx", buf.getvalue())


def png_bytes() -> tuple[str, bytes]:
    from PIL import Image

    buf = io.BytesIO()
    Image.new("RGB", (48, 24), (200, 120, 40)).save(buf, format="PNG")
    return ("image.png", buf.getvalue())
